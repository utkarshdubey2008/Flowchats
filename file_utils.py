from PIL import Image
from pptx import Presentation
from pptx.util import Inches
import os

def merge_images_to_pdf(image_files, temp_dir):
    imgs = [Image.open(img).convert("RGB") for img in image_files]
    pdf_path = os.path.join(temp_dir, "presentation.pdf")
    imgs[0].save(pdf_path, save_all=True, append_images=imgs[1:])
    return pdf_path

def merge_images_to_ppt(image_files, temp_dir):
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    for img_file in image_files:
        slide = prs.slides.add_slide(blank_slide_layout)
        slide.shapes.add_picture(img_file, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
    ppt_path = os.path.join(temp_dir, "presentation.pptx")
    prs.save(ppt_path)
    return ppt_path
